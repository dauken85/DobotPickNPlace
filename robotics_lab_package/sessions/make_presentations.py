"""
Generate PowerPoint presentations for each lab session.
Styled according to the Högskolan Väst Graphic Manual (2025).

Colour palette (from manual):
  HV Blue      #1380a4  — primary accent, logo colour, header bars
  HV Dark Navy #003b5b  — title slide backgrounds, section dividers
  HV Light     #e4f1f8  — light slide backgrounds
  Near-black   #202020  — body text on light backgrounds
  White        #ffffff  — text on coloured backgrounds

Typeface: Arial Bold / Arial Regular
  (internal standard substitute for Scout, per the manual's Internt bruk rule)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ── Högskolan Väst colour palette ────────────────────────────
HV_BLUE   = RGBColor(0x13, 0x80, 0xA4)   # #1380a4
HV_NAVY   = RGBColor(0x00, 0x3B, 0x5B)   # #003b5b
HV_LIGHT  = RGBColor(0xE4, 0xF1, 0xF8)   # #e4f1f8
HV_BLACK  = RGBColor(0x20, 0x20, 0x20)   # #202020
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
HV_MUTED  = RGBColor(0x44, 0x5F, 0x72)   # footer text


# ── Low-level helpers ─────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h,
        size=16, bold=False, italic=False,
        color=HV_BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = 'Arial'
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def bullets(slide, items, l, t, w, h, size=16, color=HV_BLACK):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"\u2013  {item}" if item.strip() else ""
        run.font.name  = 'Arial'
        run.font.size  = Pt(size)
        run.font.color.rgb = color
    return box


# ── Shared components ─────────────────────────────────────────

def footer(slide):
    """Thin HV-blue line + small footer text — per HV avdelare/footer rule."""
    rect(slide, 0, 7.22, 13.33, 0.03, HV_BLUE)
    txt(slide, "Vision-Guided Robotic Picking Lab  \u00b7  H\u00f6gskolan V\u00e4st",
        0.45, 7.27, 12.4, 0.22,
        size=9, color=HV_MUTED, align=PP_ALIGN.LEFT)


def wordmark(slide, x=11.55, y=0.20, color=WHITE):
    """'HÖGSKOLAN VÄST' text stand-in for the logo (Arial Bold uppercase)."""
    txt(slide, "H\u00d6GSKOLAN V\u00c4ST",
        x, y, 1.6, 0.35,
        size=8, bold=True, color=color, align=PP_ALIGN.RIGHT)


def hv_header(slide, heading, bg_color=HV_BLUE, text_color=WHITE, logo_color=WHITE):
    """Full-width header bar — the standard HV PPT layout header."""
    rect(slide, 0, 0, 13.33, 1.25, bg_color)
    txt(slide, heading, 0.45, 0.22, 12.0, 0.82,
        size=28, bold=True, color=text_color)
    wordmark(slide, x=11.55, y=0.42, color=logo_color)


# ── Slide builders ─────────────────────────────────────────────

def title_slide(prs, session_num, title, subtitle, duration):
    """
    Dark navy background, HV-blue top accent bar, large white title.
    Matches the dark-background 'Titel på presentationen' layout in the HV template.
    """
    slide = blank_slide(prs)
    fill_bg(slide, HV_NAVY)

    # Thin HV-blue top accent bar
    rect(slide, 0, 0, 13.33, 0.1, HV_BLUE)

    # HÖGSKOLAN VÄST wordmark top-right
    wordmark(slide, x=11.55, y=0.18, color=WHITE)

    # Session badge
    txt(slide, f"SESSION  {session_num}",
        0.55, 0.55, 4, 0.38,
        size=11, bold=True, color=HV_BLUE)

    # Thin horizontal rule below badge (HV avdelare)
    rect(slide, 0.55, 1.00, 4.5, 0.015, HV_BLUE)

    # Main title — large, white, bold Arial
    txt(slide, title,
        0.55, 1.15, 12.2, 2.5,
        size=42, bold=True, color=WHITE)

    # Subtitle — lighter weight, muted blue
    txt(slide, subtitle,
        0.55, 3.75, 10.5, 0.65,
        size=17, color=RGBColor(0xA0, 0xCC, 0xE0))

    # Duration pill — HV-blue filled box
    rect(slide, 0.55, 4.55, 2.5, 0.40, HV_BLUE)
    txt(slide, f"\u23f1  {duration}",
        0.6, 4.57, 2.4, 0.37,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Tagline bottom-left (required by HV manual in all communication)
    txt(slide, "Tillsammans f\u00f6r\u00e4ndrar vi.",
        0.55, 6.85, 6, 0.33,
        size=10, italic=True,
        color=RGBColor(0x70, 0xA8, 0xC0))

    # Closing HV-blue rule
    rect(slide, 0, 7.35, 13.33, 0.03, HV_BLUE)


def section_slide(prs, section_title):
    """
    Dark navy section divider slide — mirrors the HV 'Avsnittsrubrik' layout.
    """
    slide = blank_slide(prs)
    fill_bg(slide, HV_NAVY)
    rect(slide, 0, 0, 13.33, 0.1, HV_BLUE)
    wordmark(slide, x=11.55, y=0.18, color=WHITE)
    txt(slide, section_title,
        0.55, 2.8, 12.2, 1.8,
        size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(slide, 2.5, 4.75, 8.33, 0.02, HV_BLUE)
    txt(slide, "Tillsammans f\u00f6r\u00e4ndrar vi.",
        0.55, 6.85, 6, 0.33,
        size=10, italic=True,
        color=RGBColor(0x70, 0xA8, 0xC0))
    rect(slide, 0, 7.35, 13.33, 0.03, HV_BLUE)


def objectives_slide(prs, objectives):
    """White background, HV-blue header bar, dash bullet list."""
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    hv_header(slide, "Learning Objectives")
    bullets(slide, objectives,
            0.55, 1.45, 12.2, 5.5,
            size=17, color=HV_BLACK)
    footer(slide)


def agenda_slide(prs, parts):
    """
    Alternating HV-light / white rows on white background.
    Time in HV-blue bold, thin vertical separator, title in near-black.
    """
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    hv_header(slide, "Session Agenda", bg_color=HV_NAVY)

    row_h = 0.60
    y0    = 1.35
    for i, (time_lbl, part_title) in enumerate(parts):
        y  = y0 + i * row_h
        bg = HV_LIGHT if i % 2 == 0 else WHITE
        rect(slide, 0.35, y, 12.6, row_h - 0.06, bg)
        txt(slide, time_lbl,
            0.48, y + 0.11, 1.5, 0.38,
            size=11, bold=True, color=HV_BLUE)
        # thin vertical separator
        rect(slide, 2.12, y + 0.09, 0.018, 0.36, HV_BLUE)
        txt(slide, part_title,
            2.25, y + 0.11, 10.4, 0.38,
            size=13, color=HV_BLACK)

    footer(slide)


def content_slide(prs, heading, items):
    """
    Standard HV content slide: white background, HV-blue header, dash bullets.
    """
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    hv_header(slide, heading)
    bullets(slide, items,
            0.55, 1.45, 12.2, 5.5,
            size=16, color=HV_BLACK)
    footer(slide)
    return slide


def two_col_slide(prs, heading, left_items, right_items,
                  left_title="", right_title=""):
    """
    Two-column layout: white bg, HV-navy header bar,
    HV-light column boxes — per the HV informationsrutor element.
    """
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    hv_header(slide, heading, bg_color=HV_NAVY)

    # Left column box
    rect(slide, 0.35, 1.35, 6.05, 5.68, HV_LIGHT)
    if left_title:
        txt(slide, left_title,
            0.55, 1.42, 5.7, 0.44,
            size=13, bold=True, color=HV_NAVY)
        rect(slide, 0.55, 1.87, 5.7, 0.018, HV_BLUE)
    bullets(slide, left_items,
            0.55, 1.95 if left_title else 1.45,
            5.7, 4.8, size=13, color=HV_BLACK)

    # Right column box
    rect(slide, 6.93, 1.35, 6.05, 5.68, HV_LIGHT)
    if right_title:
        txt(slide, right_title,
            7.10, 1.42, 5.7, 0.44,
            size=13, bold=True, color=HV_NAVY)
        rect(slide, 7.10, 1.87, 5.7, 0.018, HV_BLUE)
    bullets(slide, right_items,
            7.10, 1.95 if right_title else 1.45,
            5.7, 4.8, size=13, color=HV_BLACK)

    footer(slide)


def key_takeaway_slide(prs, quote, session_num):
    """
    HV-blue full-bleed background key-message slide.
    Large white quote, tagline, session label.
    Mirrors the HV dark section intro slide style.
    """
    slide = blank_slide(prs)
    fill_bg(slide, HV_BLUE)
    rect(slide, 0, 0, 13.33, 1.15, HV_NAVY)
    wordmark(slide, x=11.55, y=0.38, color=WHITE)
    txt(slide, "Key Takeaway",
        0.5, 0.2, 12, 0.78,
        size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # HV avdelare rule
    rect(slide, 1.0, 1.28, 11.33, 0.02, WHITE)

    txt(slide, f'\u201c{quote}\u201d',
        0.9, 1.55, 11.5, 4.7,
        size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(slide, "Tillsammans f\u00f6r\u00e4ndrar vi.",
        0.55, 6.85, 6, 0.33,
        size=10, italic=True,
        color=RGBColor(0xD0, 0xEC, 0xF8))
    txt(slide, f"End of Session {session_num}",
        9.5, 6.85, 3.5, 0.33,
        size=10, color=RGBColor(0xD0, 0xEC, 0xF8),
        align=PP_ALIGN.RIGHT)
    rect(slide, 0, 7.35, 13.33, 0.03, HV_NAVY)


# ════════════════════════════════════════════════════════════════
#  SESSION 1
# ════════════════════════════════════════════════════════════════

def build_session1():
    prs = new_prs()

    title_slide(prs, 1,
                "System Overview, Safety\nand Image Capture",
                "Understanding the pipeline and building your first dataset",
                "~3 hours")

    objectives_slide(prs, [
        "Understand the full vision-to-pick pipeline",
        "Learn robot safety procedures and emergency stop",
        "Set up the camera station with consistent lighting",
        "Capture a high-quality image dataset for model training (30\u201350+ images)",
        "Understand what makes a good vs bad training image",
    ])

    agenda_slide(prs, [
        ("00:00", "Welcome & Pipeline Overview"),
        ("00:30", "Teacher Live Demo"),
        ("00:45", "Safety Briefing \u2014 every student must acknowledge"),
        ("01:00", "Station Setup \u2014 robot, camera, lighting"),
        ("01:30", "Break"),
        ("01:45", "Image Capture \u2014 varied arrangements"),
        ("02:30", "Dataset Review as a group"),
        ("03:00", "Wrap-up & preview of Session 2"),
    ])

    content_slide(prs, "The Vision-to-Pick Pipeline", [
        "\u2460  Camera captures an image of the workspace",
        "\u2461  Segmentation model detects objects and outlines their shapes",
        "\u2462  Calibration converts pixel positions \u2192 physical mm",
        "\u2463  Robot arm moves to the computed pick point and grabs the object",
        "",
        "Each session covers one or two stages of this pipeline.",
        "By Session 4 you will run the full loop automatically.",
    ])

    content_slide(prs, "Safety Rules \u2014 Read Before You Touch the Robot", [
        "Know where the Emergency Stop is \u2014 test it before starting",
        "Keep hands clear of the robot workspace during operation",
        "Never change robot speed or workspace limits without teacher approval",
        "Always tell your partner before starting any robot motion",
        "Report any unusual behaviour (jerky moves, error lights) immediately",
        "Teacher must approve your first automated pick sequence",
        "",
        "\u26a0  Each student must verbally acknowledge these rules before proceeding.",
    ])

    two_col_slide(prs,
                  "Station Setup Checklist",
                  [
                      "Robot powered on and homed",
                      "Basler camera connected and test image captured",
                      "Camera position fixed \u2014 top-down, centred on workspace",
                      "Lighting: diffuse, no strong shadows",
                      "Clean, uniform workspace background",
                      "Objects placed within camera field of view",
                  ],
                  [
                      "Run capture_images.py once",
                      "Verify image is sharp and well-exposed",
                      "Check: objects are fully visible, not cropped",
                      "Check: no mirror reflections from shiny surfaces",
                      "If something looks off \u2014 tell the teacher now",
                  ],
                  left_title="Hardware", right_title="Software / Verify")

    content_slide(prs, "Image Capture \u2014 Tips for a Good Dataset", [
        "Capture at least 30\u201350 images per group",
        "Vary arrangements: single objects, multiple objects, different positions",
        "Include edge cases: objects touching, near the border, at angles",
        "Keep lighting consistent across all captures",
        "Delete blurry, over-exposed, or poorly-lit images immediately",
        "",
        "Rule of thumb: if you can\u2019t clearly identify the object, the model can\u2019t either.",
    ])

    content_slide(prs, "What Makes a Good Training Dataset?", [
        "Diversity \u2014 many positions, orientations, and arrangements",
        "Balance \u2014 roughly equal images of each object class",
        "Quality \u2014 sharp, well-lit, consistent background",
        "Volume \u2014 more labelled data \u2192 better model (diminishing returns)",
        "Relevance \u2014 matches real operating conditions",
        "",
        "Bad data \u2192 bad model, no matter how good the algorithm.",
        "Focus on eliminating shadows and background clutter first.",
    ])

    key_takeaway_slide(prs,
        "The quality of your dataset directly determines the quality of your model. "
        "Garbage in, garbage out.",
        1)

    out = os.path.join(OUT_DIR, "Session_1_Capture.pptx")
    prs.save(out)
    print(f"Saved: {out}")


# ════════════════════════════════════════════════════════════════
#  SESSION 2
# ════════════════════════════════════════════════════════════════

def build_session2():
    prs = new_prs()

    title_slide(prs, 2,
                "Annotation and Model Training",
                "Labelling objects and teaching the model to see",
                "~3 hours")

    objectives_slide(prs, [
        "Understand what instance segmentation annotation means",
        "Label objects with accurate polygon masks in Roboflow",
        "Train a segmentation model on your own dataset",
        "Evaluate model quality using precision, recall, and mAP",
        "Identify what to fix when the model performs poorly",
    ])

    agenda_slide(prs, [
        ("00:00", "Recap of Session 1 \u2014 why dataset quality matters"),
        ("00:15", "Intro to Roboflow \u2014 upload and organise images"),
        ("00:45", "Teacher annotation demo"),
        ("01:00", "Students annotate images"),
        ("01:45", "Break"),
        ("02:00", "Train/validation split \u2014 explained"),
        ("02:15", "Start model training (runs ~15\u201330 min)"),
        ("02:30", "Evaluate results \u2014 metrics and test predictions"),
        ("03:00", "Group presentations + wrap-up"),
    ])

    two_col_slide(prs,
                  "Three Levels of Vision",
                  [
                      "Classification",
                      "  \u2192 What is this image?",
                      "  \u2192 Single label for the whole image",
                      "",
                      "Object Detection",
                      "  \u2192 Where are the objects?",
                      "  \u2192 Bounding box around each object",
                  ],
                  [
                      "Instance Segmentation  \u2190 we use this",
                      "  \u2192 Where is each object exactly?",
                      "  \u2192 Pixel-level polygon mask per object",
                      "  \u2192 Needed for accurate pick-point estimation",
                      "",
                      "More precise mask \u2192 better centroid \u2192 better pick.",
                  ],
                  left_title="Simpler approaches", right_title="What we need")

    content_slide(prs, "Annotation Best Practices", [
        "Use the smart select / polygon tool in Roboflow",
        "Trace tightly around the actual object boundary",
        "Assign the correct class name \u2014 consistent across all images",
        "Class names: square, rectangle, circle, triangle (lowercase, no spaces)",
        "A loose annotation teaches the model the background, not the object",
        "",
        "Good annotation:  tight, complete outline of the object",
        "Bad annotation:   box drawn loosely around the object",
        "",
        "Teacher will review and provide feedback during this time.",
    ])

    content_slide(prs, "Train / Validation / Test Split", [
        "Training set   (~70\u201380%)  \u2014 images the model learns from",
        "Validation set (~10\u201315%)  \u2014 used to tune hyperparameters during training",
        "Test set       (~10\u201315%)  \u2014 held-out images for final evaluation",
        "",
        "Key rule: NEVER test on the same images you trained on",
        "          (the model would just memorise the answers)",
        "",
        "Roboflow handles the split automatically \u2014 accept the defaults.",
    ])

    content_slide(prs, "Understanding Model Metrics", [
        "mAP  (mean Average Precision) \u2014 overall quality score",
        "      > 0.85  Excellent      0.70\u20130.85  Good      < 0.70  Needs work",
        "",
        "Precision \u2014 When the model says \u2018circle\u2019, how often is it right?",
        "            High precision = few false positives",
        "",
        "Recall    \u2014 Of all the circles in the image, how many did it find?",
        "            High recall = few missed objects",
        "",
        "For pick-and-place: high recall matters most (missing an object is bad).",
    ])

    content_slide(prs, "Evaluating Your Model \u2014 What to Look For", [
        "Run predictions on the test set in Roboflow",
        "Check each image for:",
        "  \u2713  Correct detections with tight mask boundaries",
        "  \u2717  False positives (object detected where there is none)",
        "  \u2717  Missed objects (object not detected at all)",
        "  \u2717  Wrong class label assigned",
        "  \u2717  Mask extending far beyond the object",
        "",
        "If precision is low \u2192 add more varied training images",
        "If recall is low   \u2192 add more examples of the missed class",
    ])

    key_takeaway_slide(prs,
        "A model is only as good as its annotations. "
        "Sloppy labels make a sloppy model \u2014 "
        "tighter, more accurate masks produce better segmentation.",
        2)

    out = os.path.join(OUT_DIR, "Session_2_Annotation.pptx")
    prs.save(out)
    print(f"Saved: {out}")


# ════════════════════════════════════════════════════════════════
#  SESSION 3
# ════════════════════════════════════════════════════════════════

def build_session3():
    prs = new_prs()

    title_slide(prs, 3,
                "Inference and ArUco Calibration",
                "From pixels to millimetres \u2014 bridging the gap between camera and robot",
                "~3 hours")

    objectives_slide(prs, [
        "Run a trained segmentation model on images and interpret the output",
        "Understand the pixel-to-mm conversion problem",
        "Perform ArUco-based calibration to compute mm_per_pixel",
        "Convert detected object positions from pixels to physical millimetres",
        "Verify calibration accuracy against a ruler measurement",
    ])

    agenda_slide(prs, [
        ("00:00", "Recap \u2014 we have a model, now what?"),
        ("00:15", "Inference on saved images \u2014 confidence threshold exercise"),
        ("00:45", "The calibration problem \u2014 pixels vs mm"),
        ("01:00", "Break"),
        ("01:15", "ArUco calibration exercise"),
        ("02:00", "Apply calibration to inference results"),
        ("02:30", "Live camera inference"),
        ("02:45", "Group results + wrap-up"),
    ])

    content_slide(prs, "Running Inference \u2014 What the Model Returns", [
        "Input:  an image (JPEG/PNG) sent to the Roboflow API",
        "Output (JSON) for each detected object:",
        "  \u2022 class       \u2014 the class label (e.g. \u2018circle\u2019)",
        "  \u2022 confidence  \u2014 model certainty 0.0\u20131.0",
        "  \u2022 points      \u2014 list of (x, y) polygon vertices in pixels",
        "",
        "From the polygon we compute the centroid (cx, cy) using image moments.",
        "",
        "Exercise: change confidence threshold 0.3 \u2192 0.5 \u2192 0.8",
        "         Observe how detection count changes.",
    ])

    content_slide(prs, "The Calibration Problem", [
        "The model outputs positions in pixels  (e.g. cx=847 px, cy=512 px)",
        "The robot needs positions in millimetres  (e.g. x=215 mm, y=130 mm)",
        "",
        "We need:    mm_per_pixel  =  physical_distance_mm / distance_pixels",
        "",
        "Assumptions that must hold:",
        "  \u2022 Camera is fixed between capture and pick",
        "  \u2022 Workspace is flat \u2014 all objects on the same plane",
        "  \u2022 Camera is mounted top-down (no perspective distortion)",
        "",
        "If any assumption is violated, accuracy will degrade.",
    ])

    content_slide(prs, "ArUco Markers \u2014 How Calibration Works", [
        "An ArUco marker is a square fiducial with a unique ID pattern",
        "We print one at a known physical size (e.g. 50 mm \u00d7 50 mm)",
        "",
        "Steps:",
        "  1. Place the marker flat in the workspace",
        "  2. Capture an image",
        "  3. Detect the 4 corner points of the marker in pixels",
        "  4. Measure the side length in pixels  (C0\u2192C1 distance)",
        "  5. Compute:  mm_per_pixel = 50 mm / side_px",
        "",
        "Corner C0 of the marker = work object origin (0, 0)",
        "All object positions are reported relative to C0.",
    ])

    two_col_slide(prs,
                  "Student Exercise \u2014 ArUco Calibration",
                  [
                      "\u2460 Measure the printed marker with a ruler",
                      "   \u2192 Record the side length in mm",
                      "",
                      "\u2461 Enter the measurement in the notebook",
                      "   MARKER_SIZE_MM = ___",
                      "",
                      "\u2462 Run the calibration cell",
                      "   \u2192 Note mm_per_pixel value",
                      "",
                      "\u2463 Record your result in the group log",
                  ],
                  [
                      "Verification step:",
                      "",
                      "Place an object at a known distance from C0",
                      "Measure with a ruler  \u2192  X mm",
                      "",
                      "Run verification cell",
                      "Compare computed vs measured distance",
                      "",
                      "Goal: error < 5%",
                      "If > 10% \u2192 recheck marker size or camera tilt",
                  ],
                  left_title="Calibration", right_title="Verification")

    content_slide(prs, "Apply Calibration to Inference Results", [
        "After calibration you have:  mm_per_pixel  and  c0_px (origin in pixels)",
        "",
        "For each detected object centroid (px_x, px_y):",
        "",
        "  wobj_x = (px_x \u2212 c0_px[0])  \u00d7  mm_per_pixel",
        "  wobj_y = (px_y \u2212 c0_px[1])  \u00d7  mm_per_pixel",
        "",
        "These coordinates are directly in the robot work object frame.",
        "No manual offset is needed when the work object is defined at C0.",
        "",
        "Activate with:  User(N)  on the robot before sending move commands.",
    ])

    key_takeaway_slide(prs,
        "One printed marker, one ruler measurement, and one division "
        "turns every pixel in your image into a real-world coordinate "
        "the robot understands.",
        3)

    out = os.path.join(OUT_DIR, "Session_3_Inference_Calibration.pptx")
    prs.save(out)
    print(f"Saved: {out}")


# ════════════════════════════════════════════════════════════════
#  SESSION 4
# ════════════════════════════════════════════════════════════════

def build_session4():
    prs = new_prs()

    title_slide(prs, 4,
                "Robot-Guided Pick and Place",
                "Connecting vision to motion \u2014 the full autonomous pipeline",
                "~3 hours")

    objectives_slide(prs, [
        "Connect the complete vision pipeline to the Dobot E6 robot arm",
        "Execute supervised and autonomous pick-and-place sequences",
        "Use the ArUco work object frame for accurate positioning",
        "Log pick results and compute a success rate",
        "Analyse failures and propose improvements",
    ])

    agenda_slide(prs, [
        ("00:00", "Recap + Safety Review \u2014 mandatory before robot use"),
        ("00:15", "Teacher verifies every station individually"),
        ("00:45", "First supervised picks \u2014 teacher present"),
        ("01:15", "Break"),
        ("01:30", "Independent picking \u2014 increasing difficulty"),
        ("02:00", "Mixed clutter challenge \u2014 5+ objects"),
        ("02:30", "Final group demo (2 min each)"),
        ("02:45", "Discussion, wrap-up, and reflection"),
    ])

    content_slide(prs, "Safety Review \u2014 Robot Section", [
        "\u26a0  Review before EVERY robot session",
        "",
        "Emergency stop location confirmed",
        "Hands and tools clear of workspace before any move",
        "Robot speed set to a safe value (\u226430%)",
        "Teacher has approved the first pick at your station",
        "Drop zone is defined and clear",
        "",
        "If the robot moves unexpectedly \u2192 press E-stop immediately",
        "Do not try to catch a falling object with your hands",
        "",
        "Teacher must countersign your station checklist before you begin.",
    ])

    content_slide(prs, "The Full Pipeline \u2014 All Four Steps Together", [
        "\u2460 Capture       \u2014 Take a fresh image from the Basler camera",
        "\u2461 Segment       \u2014 Send image to Roboflow \u2192 receive mask polygons",
        "\u2462 Compute       \u2014 Convert centroids to work object coordinates",
        "                  wobj_x = (px_x \u2212 c0_px[0]) \u00d7 mm_per_pixel",
        "\u2463 Pick          \u2014 Send MovL command to Dobot E6",
        "                  Approach \u2192 Descend \u2192 Suction ON \u2192 Lift \u2192 Drop \u2192 OFF",
        "",
        "The robot uses User coordinate frame N \u2014 C0 of the ArUco marker is (0, 0, 0).",
    ])

    two_col_slide(prs,
                  "First Picks \u2014 Step by Step",
                  [
                      "Place ONE object in a clear area",
                      "Run the full pipeline cell by cell",
                      "Observe the pick point on the visualisation",
                      "Teacher approves the target before robot moves",
                      "Record the outcome in your pick log",
                      "",
                      "Adjust X/Y/Z offsets if needed",
                      "Repeat for 3\u20135 single-object picks",
                  ],
                  [
                      "Pick Log format:",
                      "",
                      "Pick# | Class  | Conf | Position | Result",
                      "1     | circle | 0.92 | (45, 30) | \u2713",
                      "2     | square | 0.78 | (80, 55) | \u2717 off 8mm",
                      "",
                      "Track every attempt.",
                      "You will present your success rate at the end.",
                  ],
                  left_title="Procedure", right_title="Recording Results")

    content_slide(prs, "Independent Picking \u2014 Increasing Complexity", [
        "Level 1 \u2014 Single object, centred in workspace",
        "Level 2 \u2014 Two objects: pick the one with higher confidence",
        "Level 3 \u2014 Three objects: pick in order (e.g. largest area first)",
        "Level 4 \u2014 Objects near the workspace boundary",
        "",
        "For each level: run the pipeline \u2192 log the result \u2192 analyse failures",
        "",
        "If a pick fails, diagnose before continuing:",
        "  \u2022 Was the detection correct? (check the visualisation)",
        "  \u2022 Was the computed coordinate plausible?",
        "  \u2022 Did the suction cup miss the object centroid?",
    ])

    content_slide(prs, "Mixed Clutter Challenge", [
        "Arrange 5 or more objects \u2014 some touching, some overlapping",
        "Mix all object classes together",
        "Run the full pipeline and attempt to pick every object",
        "",
        "Discussion questions after your attempt:",
        "  \u2022 Which objects were picked successfully?",
        "  \u2022 Which were missed or misidentified?",
        "  \u2022 Did overlapping objects cause problems?",
        "  \u2022 Would a different pick order have worked better?",
        "  \u2022 What would need to change for a real factory setup?",
    ])

    content_slide(prs, "Wrap-up Discussion \u2014 Key Questions", [
        "What was the single biggest source of error at your station?",
        "How did your mm_per_pixel calibration compare to other groups?",
        "What would you change to improve pick reliability to > 90%?",
        "",
        "Looking further ahead:",
        "  \u2022 How would you handle 3D objects (not flat on a surface)?",
        "  \u2022 What if a new object type appears that was not in training?",
        "  \u2022 How would a real factory system handle lighting changes?",
        "  \u2022 What safety certifications would a production robot need?",
    ])

    key_takeaway_slide(prs,
        "Vision-guided robotics is not magic \u2014 it is a chain of steps, "
        "each introducing a small error. "
        "Understanding where errors come from is the first step to eliminating them.",
        4)

    out = os.path.join(OUT_DIR, "Session_4_Pick_and_Place.pptx")
    prs.save(out)
    print(f"Saved: {out}")


# ── Main ──────────────────────────────────────────────────────
if __name__ == "__main__":
    build_session1()
    build_session2()
    build_session3()
    build_session4()
    print("\nAll presentations generated successfully.")
